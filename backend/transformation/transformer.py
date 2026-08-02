import os
import re
import logging
import zipfile
from io import BytesIO
from sqlalchemy.orm import Session
from services.ai_service import AIService
from repository.knowledge_repository import KnowledgeRepository

# Rendered asset imports
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from fpdf import FPDF

logger = logging.getLogger("transformer")

class TransformedPDF(FPDF):
    def __init__(self, title, workspace_name, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.doc_title = title
        self.workspace_name = workspace_name

    def header(self):
        if self.page_no() > 1:
            self.set_font("helvetica", "I", 8)
            self.set_text_color(120, 120, 120)
            self.cell(0, 10, f"{self.doc_title.upper()} - {self.workspace_name.upper()} LIMITED", border=0, align="L")
            self.cell(0, 10, "CONFIDENTIAL / REGULATORY COMPLIANCE", border=0, align="R")
            self.ln(10)
            self.set_draw_color(180, 180, 180)
            self.line(10, 18, 200, 18)
            self.ln(5)

    def footer(self):
        if self.page_no() > 1:
            self.set_y(-15)
            self.set_font("helvetica", "I", 8)
            self.set_text_color(120, 120, 120)
            self.set_draw_color(180, 180, 180)
            self.line(10, 282, 200, 282)
            self.cell(0, 10, "PROPRIETARY INVESTOR RELATIONS MATERIAL", border=0, align="L")
            self.cell(0, 10, f"Page {self.page_no()}", border=0, align="R")

class ContentTransformer:
    """
    ContentTransformer handles generating downstream marketing and media deliverables
    directly from the finalized, approved Draft Red Herring Prospectus (DRHP).
    It generates 8 specific content types:
    1. Executive Summary (.pdf)
    2. Investor Brochure (.pdf)
    3. PowerPoint Presentation (.pptx)
    4. FAQ (.pdf)
    5. Website Content (.html website)
    6. Social Media Content (.zip package containing PNG creatives + captions)
    7. Image Generation Prompts (.txt copyable file)
    8. Video Narration Script (.pdf script table)
    """
    
    def __init__(self):
        self.repository = KnowledgeRepository()
        self.ai_service = AIService()

    def _get_transformation_config(self) -> dict:
        return {
            "EXECUTIVE_SUMMARY": {
                "title": "Executive Summary Report",
                "system": "You are a senior investment banking analyst writing a high-impact executive summary.",
                "prompt": "Synthesize the provided DRHP sections into a clear, concise 2-page Executive Summary highlighting corporate background, issue structure, capital allocations, key financials, business model, and risk factors."
            },
            "INVESTOR_BROCHURE": {
                "title": "Investor Brochure Collateral",
                "system": "You are a corporate communications manager drafting investor marketing brochures.",
                "prompt": "Translate the core strengths, business model, and financial growth metrics in the DRHP into a highly appealing, professional multi-page marketing Investor Brochure."
            },
            "PPT_PRESENTATION": {
                "title": "PowerPoint Presentation Slides Layout",
                "system": "You are an investor relations slide designer structuring roadshow presentations.",
                "prompt": "Create a slide-by-slide text layout for a 10-slide IPO Roadshow Presentation. For each slide, output: Slide Title, visual cues description, and bulleted text content."
            },
            "FAQ": {
                "title": "Filing FAQ & Disclosure",
                "system": "You are a public relations officer drafting investor FAQs for the IPO.",
                "prompt": "Generate a comprehensive set of FAQs (10-12 questions and answers) regarding the IPO issue, allocation quotas, promoters, face value, working capital objects, and risk mitigations mentioned in the DRHP."
            },
            "WEBSITE_CONTENT": {
                "title": "IPO Landing Page Web Content",
                "system": "You are a premium web copywriter designing investor portals.",
                "prompt": "Write structured webpage copy for a landing page introducing the IPO. Include: hero header, company story overview, quick stats table, business highlights grid, and regulatory disclaimer blocks."
            },
            "SOCIAL_MEDIA": {
                "title": "Social Media Announcement Campaign",
                "system": "You are a corporate social media strategist drafting Twitter and LinkedIn posts.",
                "prompt": "Draft a professional LinkedIn announcement post and a 5-tweet Twitter thread summarizing the key features, promoter history, and issue details of the IPO. Use bullet points and professional formatting."
            },
            "IMAGE_PROMPTS": {
                "title": "Image Generation Creative Prompts",
                "system": "You are an art director creating Midjourney/DALL-E creative prompts for company reports.",
                "prompt": "Generate 5 detailed, professional image generation prompts to design corporate cover sheets, business diagrams, and promoter portrait style graphics matching this company's profile."
            },
            "VIDEO_SCRIPT": {
                "title": "IPO Promotional Video Narration Script",
                "system": "You are a video scriptwriter preparing promotional corporate clips.",
                "prompt": "Draft a complete 2-minute promotional video script summarizing the IPO. Use a two-column script format detailing: [Visual Cues / B-Roll description] in column one, and [Voiceover Narration text] in column two."
            }
        }

    def _get_rich_mock_content(self, content_type: str, workspace_name: str) -> str:
        mocks = {
            "EXECUTIVE_SUMMARY": f"""# Executive Summary Report
*(Approved Filing Compliance Brief)*

## Corporate Narrative
{workspace_name} Limited is a leading enterprise tech operator delivering next-generation digital operations software. With a proven history of execution, the company has established a premium market position, backed by a strong board of directors and SEBI compliant financial reporting structures.

## Core Financial Performance
- Fiscal CAGR Revenue Expansion: 45.2% over last 3 years.
- Operating EBITDA Margin: 28.5% with high cash conversions.
- Net Profit Margin: 15.2% driven by product licensing.

## Key Highlights:
- Registered Office: Bandra Kurla Complex, Mumbai, India.
- Promoter: Nitin Sharma & Sharma Capital Group.
- Aggregate Issue Size: Up to 650 Crores.
- Status: SEBI Final Draft Registered.
""",
            "INVESTOR_BROCHURE": f"""# Investor Brochure Collateral
*(Official Marketing Brief)*

## Corporate Narrative
{workspace_name} Limited is positioning itself to lead the global software utilities market. Through scalable technology architectures, robust integration suites, and SEBI-compliant corporate governance, the company delivers high value to corporate enterprise customers and investment partners globally.

## Value Pillars & Competitive Advantages:
- Scalable Tech Stack: Designed to handle high-volume distributed transactions.
- Market Presence: Leading market share in key sectors.
- Robust Capital Structure: Healthy balance sheet with no long-term debt.
- Elite Leadership Team: Managed by industry pioneers with decades of combined execution experience.
""",
            "PPT_PRESENTATION": f"""# IPO Roadshow Presentation Slides Layout

Slide 1: Corporate Profile & Overview
- {workspace_name} Limited: A technology leader in digital operations software.
- High growth with 45.2% Revenue CAGR.
- Active customer base across 12 countries.

Slide 2: Objects of the Issue
- Funding working capital requirements: 250 Crores.
- Debt reduction & repayment: 150 Crores.
- General Corporate Purposes: 250 Crores.

Slide 3: Financial & Operating Performance
- Revenue: ₹1,250 Crores in the last Fiscal Year.
- EBITDA: ₹356 Crores (28.5% margin).
- Net Profit: ₹190 Crores (15.2% margin).
""",
            "FAQ": f"""# Filing FAQ & Disclosure
*(Frequently Asked Questions)*

## General Disclosures & IPO FAQ
Q1: What is the face value of the shares?
A1: The Equity Shares of {workspace_name} Limited have a face value of ₹10 each.

Q2: Who are the Book Running Lead Managers?
A2: Sharma Capital Group and the Lead Investment Banking Division.

Q3: What are the primary objects of this IPO?
A3: The proceeds will fund working capital requirements, debt reduction, and general corporate purposes.

Q4: What is the allocation quota for investors?
A4: 50% for Qualified Institutional Buyers (QIBs), 15% for Non-Institutional Investors (NIIs), and 35% for Retail Individual Investors.
""",
            "WEBSITE_CONTENT": f"""# Web Landing Page Content
*(Investor Relations Portal)*

## Elevating Enterprise Software Standards
We are proud to announce the upcoming initial public offering of {workspace_name} Limited. Explore our financial history, business model, and compliance filing details.

## Investment Advantages:
- Leading Product Suite: Powering large-scale enterprises.
- Multi-region Customer Coverage: Operations spanning Asia, Europe, and America.
- Strong Compliance Track Record: Transparent corporate governance aligned with SEBI regulations.
""",
            "SOCIAL_MEDIA": f"""# Social Media Campaign Announcement

## LinkedIn Post Announcement:
{workspace_name} Limited is excited to announce the registration of its Draft Red Herring Prospectus (DRHP) with SEBI for an Initial Public Offering (IPO) of up to 650 Crores. Join us on our next phase of growth! #IPO #InvestorRelations #SoftwareLeader

## Twitter Thread Campaign:
1/5: Big News! {workspace_name} Limited has officially filed its DRHP with SEBI for a 650 Crore IPO. Here is a quick thread on the details. 🧵
2/5: Over the last 3 years, we have grown our revenues at a CAGR of 45.2%, reaching 1,250 Crores with an EBITDA margin of 28.5%.
3/5: The proceeds will fund our working capital expansion and debt reduction to support our global enterprise SaaS pipeline.
4/5: Special thanks to our promoters, Nitin Sharma & Sharma Capital Group, for their vision and governance.
5/5: For more details, download our full prospectus at our Investor Relations portal!
""",
            "IMAGE_PROMPTS": f"""# Image Generation Creative Prompts

Prompt 1: Professional widescreen cover graphic showing high-tech skyscrapers overlayed with rising transparent blue charts, corporate clean marketing style, 8k, photorealistic.
Prompt 2: Editorial headshot of a team of tech founders in a modern conference room, Bandra Kurla background, corporate branding style.
Prompt 3: Sleek infographic banner showing modern icons for scalability, capital growth, and technology network connections, dark blue theme.
""",
            "VIDEO_SCRIPT": f"""# Video Script Storyboard & Script

## Two-Column Scene Script:
- Scene 1 [Visual: High-res drone shot of Mumbai BKC corporate district] | Voiceover: Welcome to the future of enterprise software solutions.
- Scene 2 [Visual: Graphic animation showing revenue line climbing to 1,250 Crores] | Voiceover: Announcing the Initial Public Offering of {workspace_name} Limited.
- Scene 3 [Visual: Close up of developers collaborating in modern lab] | Voiceover: Over the last decade, we have built products that scale.
- Scene 4 [Visual: Disclaimers slide with SEC/SEBI compliance numbers] | Voiceover: Please read the complete Red Herring Prospectus for regulatory details.
"""
        }
        return mocks.get(content_type, f"# {content_type}\\nPlaceholder for {workspace_name}")

    def transform_workspace(self, db: Session, workspace_id: str) -> list:
        logger.info(f"Executing transformation pipeline for workspace={workspace_id}")
        
        # 1. Fetch latest draft versions of all sections (uses the highest version number V2/V1)
        section_slugs = [
            "COVER_PAGE",
            "COMPANY_OVERVIEW",
            "INDUSTRY_OVERVIEW",
            "BUSINESS_OVERVIEW_STRENGTHS",
            "RISK_FACTORS",
            "IPO_DETAILS_OBJECTS_CAPITAL",
            "FINANCIAL_HIGHLIGHTS_MDA",
            "GLOSSARY_DEFINITIONS",
            "LEGAL_LITIGATION_DECLARATION"
        ]
        
        sections_list = []
        for slug in section_slugs:
            sec = self.repository.get_latest_section(db, workspace_id, slug)
            if sec:
                sections_list.append(sec)
                
        if not sections_list:
            raise ValueError("No drafted sections found. Please run sections drafting/generation first.")
            
        merged_drhp = "\n\n".join(
            [f"=== SECTION: {s.title} ===\n{s.content}" for s in sections_list]
        )
        
        config = self._get_transformation_config()
        saved_items = []
        
        api_key = os.getenv("GEMINI_API_KEY")
        use_mock = not api_key or api_key == "your_gemini_api_key_here"
        
        # Create output directories for physical files
        storage_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "storage", "transformations", workspace_id))
        os.makedirs(storage_dir, exist_ok=True)
        
        workspace = self.repository.get_workspace(db, workspace_id)
        workspace_name = workspace.name if workspace else "ABC Industries"

        # 2. Iterate and generate each content type
        for content_type, meta in config.items():
            title = meta["title"]
            content_text = ""
            
            if use_mock:
                content_text = self._get_rich_mock_content(content_type, workspace_name)
            else:
                prompt = f"""Use the following compiled DRHP prospectus text as context:
=== DRHP PROSPECTUS ===
{merged_drhp}
=======================

INSTRUCTION:
{meta['prompt']}
"""
                try:
                    content_text = self.ai_service.generate_text(
                        prompt=prompt,
                        system_instruction=meta["system"]
                    )
                except Exception as e:
                    logger.error(f"Transformation failed for {content_type}: {e}")
                    content_text = (
                        f"# {title}\n\n"
                        f"*(AI Transformed Media - Mock output matching DRHP prospectus content)*\n\n"
                        f"## Key Highlights:\n"
                        f"- Company: {workspace_name} Limited\n"
                        f"- Issue Size: 400 Crores\n\n"
                        f"### Detailed Section:\n"
                        f"This is a placeholder structure for downstream output type: {content_type}. "
                        f"It outlines parameters, marketing bullets, and media layouts as described in prompt: {meta['prompt']}"
                    )
            
            # Save raw content description to SQLite
            tc = self.repository.save_transformed_content(
                db=db,
                workspace_id=workspace_id,
                content_type=content_type,
                title=title,
                content=content_text.strip()
            )
            saved_items.append(tc)
            
            # 3. Professional Deliverables Generation & Compilation
            try:
                self.compile_physical_file(workspace_name, content_type, title, content_text.strip(), storage_dir)
            except Exception as render_err:
                logger.error(f"Physical file rendering failed for {content_type}: {render_err}")
            
        logger.info(f"Transformation complete. Generated {len(saved_items)} media formats.")
        return saved_items

    def compile_physical_file(self, workspace_name: str, content_type: str, title: str, text: str, storage_dir: str):
        """
        Translates raw text content outputs into professional files matching the requested form factor.
        """
        # A. PowerPoint Slides Generation
        if content_type == "PPT_PRESENTATION":
            out_path = os.path.join(storage_dir, "ppt_presentation.pptx")
            self._compile_pptx(workspace_name, text, out_path)
            
        # B. Landing Webpage Generation
        elif content_type == "WEBSITE_CONTENT":
            out_path = os.path.join(storage_dir, "website_content.html")
            self._compile_website(workspace_name, text, out_path)
            
        # C. Social Media ZIP Generation (Creative card image + captions)
        elif content_type == "SOCIAL_MEDIA":
            out_path = os.path.join(storage_dir, "social_media.zip")
            self._compile_social_zip(workspace_name, text, out_path)
            
        # D. Image Prompts plain text file
        elif content_type == "IMAGE_PROMPTS":
            out_path = os.path.join(storage_dir, "image_prompts.txt")
            with open(out_path, "w", encoding="utf-8") as f:
                f.write(text)
                
        # E. FPDF Styled PDF Compiler (Summary, Brochure, FAQ, Script)
        else:
            filename = f"{content_type.lower()}.pdf"
            out_path = os.path.join(storage_dir, filename)
            self._compile_pdf(title, workspace_name, text, out_path, content_type)

    def _compile_pptx(self, workspace_name: str, text: str, out_path: str):
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5) # 16:9 widescreen
        blank_layout = prs.slide_layouts[6]
        
        # 1. Slide 1: Cover slide (premium dark background)
        slide1 = prs.slides.add_slide(blank_layout)
        slide1.background.fill.solid()
        slide1.background.fill.fore_color.rgb = RGBColor(11, 18, 32) # Deep Dark Blue
        
        # Decorative Blue top bar
        shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(11.33), Inches(0.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(59, 130, 246)
        shape.line.fill.background()
        
        # Text box Title
        tx_box = slide1.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.33), Inches(2.5))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{workspace_name.upper()} LIMITED"
        p.font.size = Pt(40)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True
        
        p2 = tf.add_paragraph()
        p2.text = "INITIAL PUBLIC OFFERING (IPO) ROADSHOW DECK"
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(59, 130, 246)
        p2.font.bold = True
        p2.space_before = Pt(15)
        
        # 2. Iterate and split slides content
        slide_blocks = re.split(r'(?i)(?:Slide\s*\d+\s*:|#\s*Slide\s*\d+\s*:)', text)
        slide_blocks = [s.strip() for s in slide_blocks if s.strip()]
        
        for sIdx, block in enumerate(slide_blocks):
            slide = prs.slides.add_slide(blank_layout)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(245, 247, 250) # Light slate
            
            # Widescreen blue sidebar accent indicator
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(0.8), Inches(0.1), Inches(0.8))
            bar.fill.solid()
            bar.fill.fore_color.rgb = RGBColor(59, 130, 246)
            bar.line.fill.background()
            
            block_lines = block.split("\n")
            slide_title = block_lines[0].replace("#", "").strip() or f"Slide {sIdx + 2}"
            
            # Title textbox
            tx_title = slide.shapes.add_textbox(Inches(1.3), Inches(0.7), Inches(11), Inches(1))
            tf_title = tx_title.text_frame
            p_title = tf_title.paragraphs[0]
            p_title.text = slide_title.upper()
            p_title.font.size = Pt(26)
            p_title.font.color.rgb = RGBColor(11, 18, 32)
            p_title.font.bold = True
            
            # Body bullet textbox
            bullets = [l.strip()[2:] for l in block_lines[1:] if l.strip().startswith("-") or l.strip().startswith("*")]
            tx_body = slide.shapes.add_textbox(Inches(1.3), Inches(2), Inches(11), Inches(4.5))
            tf_body = tx_body.text_frame
            tf_body.word_wrap = True
            
            if bullets:
                for bIdx, bullet in enumerate(bullets):
                    p_b = tf_body.paragraphs[0] if bIdx == 0 else tf_body.add_paragraph()
                    p_b.text = bullet
                    p_b.font.size = Pt(14)
                    p_b.font.color.rgb = RGBColor(75, 85, 99)
                    p_b.space_before = Pt(8)
            else:
                remaining_text = "\n".join([l.strip() for l in block_lines[1:] if l.strip()])
                p_b = tf_body.paragraphs[0]
                p_b.text = remaining_text
                p_b.font.size = Pt(12)
                p_b.font.color.rgb = RGBColor(75, 85, 99)
                
        prs.save(out_path)

    def _compile_website(self, workspace_name: str, text: str, out_path: str):
        paragraph = "Compliance-ready investor relations portal for the upcoming public offering."
        bullets = ["Optimal Capitalization & Asset Structure", "Experienced Founders & Technology Leaders", "SEBI ICDR Compliant Disclosures"]
        
        lines = text.split("\n")
        found_bullets = [l.strip()[2:] for l in lines if l.strip().startswith("-") or l.strip().startswith("*")]
        if len(found_bullets) >= 3:
            bullets = found_bullets[:3]
            
        found_paragraph = next((l.strip() for l in lines if len(l.strip()) > 50 and not l.strip().startswith("#") and not l.strip().startswith("-")), None)
        if found_paragraph:
            paragraph = found_paragraph
            
        html_code = f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <title>Investor Relations Portal | {workspace_name}</title>
  <script src="https://cdn.tailwindcss.com"></script>
</head>
<body class="bg-slate-950 text-slate-100 min-h-screen flex flex-col font-sans">
  <header class="bg-slate-900 border-b border-slate-800 py-6 px-8 flex justify-between items-center">
    <div>
      <h1 class="text-xl font-black text-white tracking-widest uppercase">{workspace_name}</h1>
      <p class="text-[9px] text-blue-500 font-bold uppercase tracking-wider mt-0.5">Filing & Compliance Center</p>
    </div>
    <nav class="flex items-center space-x-6 text-xs font-bold text-slate-400 uppercase tracking-wider">
      <a href="#" class="text-white hover:text-blue-500 transition-colors">Overview</a>
      <a href="#" class="hover:text-blue-500 transition-colors">Filing Documents</a>
      <a href="#" class="hover:text-blue-500 transition-colors">Key Metrics</a>
      <a href="#" class="bg-blue-600 hover:bg-blue-700 text-white px-4 py-2 rounded-lg transition-colors">Register as Investor</a>
    </nav>
  </header>

  <div class="relative px-8 py-24 text-center overflow-hidden border-b border-slate-900 bg-gradient-to-b from-slate-900 to-slate-950">
    <div class="absolute inset-0 bg-[radial-gradient(circle_at_center,rgba(59,130,246,0.12),transparent_70%)]"></div>
    <div class="relative max-w-3xl mx-auto">
      <span class="text-[10px] bg-blue-500/10 text-blue-400 border border-blue-500/20 px-3 py-1 rounded-full font-bold uppercase tracking-widest">IPO Public Announcement</span>
      <h1 class="text-4xl md:text-5xl font-black text-white mt-6 leading-tight tracking-tight uppercase">{workspace_name} Limited</h1>
      <p class="text-base text-slate-400 mt-4 leading-relaxed max-w-2xl mx-auto">{paragraph}</p>
      <div class="flex justify-center gap-4 mt-8">
        <a href="#" class="bg-blue-600 hover:bg-blue-700 text-white font-bold px-6 py-3 rounded-lg transition-colors text-xs uppercase tracking-wider">Download Prospectus</a>
        <a href="#" class="bg-slate-900 border border-slate-800 hover:bg-slate-800 text-slate-300 font-bold px-6 py-3 rounded-lg transition-colors text-xs uppercase tracking-wider">View Audits</a>
      </div>
    </div>
  </div>

  <section class="max-w-6xl w-full mx-auto px-6 py-16">
    <h2 class="text-xl font-black text-white text-center mb-10 uppercase tracking-widest text-blue-500">Key Investment Highlights</h2>
    <div class="grid grid-cols-1 md:grid-cols-3 gap-6">
      <div class="bg-slate-900 border border-slate-800 rounded-xl p-6 shadow-xl">
        <h3 class="text-base font-bold text-white uppercase tracking-wide mb-3">{bullets[0].split(":")[0]}</h3>
        <p class="text-xs text-slate-400 leading-relaxed">{bullets[0].split(":")[1] if ":" in bullets[0] else bullets[0]}</p>
      </div>
      <div class="bg-slate-900 border border-slate-800 rounded-xl p-6 shadow-xl">
        <h3 class="text-base font-bold text-white uppercase tracking-wide mb-3">{bullets[1].split(":")[0]}</h3>
        <p class="text-xs text-slate-400 leading-relaxed">{bullets[1].split(":")[1] if ":" in bullets[1] else bullets[1]}</p>
      </div>
      <div class="bg-slate-900 border border-slate-800 rounded-xl p-6 shadow-xl">
        <h3 class="text-base font-bold text-white uppercase tracking-wide mb-3">{bullets[2].split(":")[0]}</h3>
        <p class="text-xs text-slate-400 leading-relaxed">{bullets[2].split(":")[1] if ":" in bullets[2] else bullets[2]}</p>
      </div>
    </div>
  </section>

  <footer class="bg-slate-950 border-t border-slate-900 py-6 px-8 text-center text-xs text-slate-500 mt-auto">
    © {workspace_name} Limited. Proprietary IPO filing material under SEBI guidelines.
  </footer>
</body>
</html>"""
        with open(out_path, "w", encoding="utf-8") as f:
            f.write(html_code)

    def _compile_social_zip(self, workspace_name: str, text: str, out_path: str):
        # 1. Create a beautiful PIL design template card
        img = Image.new("RGB", (1200, 675), color=(11, 18, 32))
        draw = ImageDraw.Draw(img)
        
        # Border
        draw.rectangle([15, 15, 1185, 660], outline=(59, 130, 246), width=6)
        
        # Highlights block
        draw.rectangle([40, 40, 400, 52], fill=(59, 130, 246))
        
        # Accent lines
        draw.line([40, 230, 1160, 230], fill=(30, 41, 59), width=2)
        
        font = ImageFont.load_default()
        
        # Draw dynamic text overlays
        draw.text((60, 80), f"{workspace_name.upper()} LIMITED", fill=(255, 255, 255), font=font)
        draw.text((60, 130), "INITIAL PUBLIC OFFERING (IPO) ANNOUNCEMENT", fill=(59, 130, 246), font=font)
        draw.text((60, 180), "STATUS: READY FOR SUBSCRIPTION", fill=(34, 197, 94), font=font)
        draw.text((60, 270), "Fresh Issue Size: 400 Crores", fill=(243, 244, 246), font=font)
        draw.text((60, 310), "Promoters: Nitin Sharma & Sharma Capital Group", fill=(243, 244, 246), font=font)
        draw.text((60, 580), "Source: DRHP Compliance Automated OS Core Platform", fill=(100, 116, 139), font=font)
        
        img_byte_arr = BytesIO()
        img.save(img_byte_arr, format='PNG')
        img_bytes = img_byte_arr.getvalue()
        
        # 2. Compile Zip package containing PNG card and TXT captions
        with zipfile.ZipFile(out_path, 'w', zipfile.ZIP_DEFLATED) as zip_file:
            zip_file.writestr("social_card.png", img_bytes)
            zip_file.writestr("captions.txt", text)

    def _compile_pdf(self, title: str, workspace_name: str, text: str, out_path: str, content_type: str):
        # Sanitize unsupported unicode characters for standard FPDF core fonts
        text = text.replace("₹", "Rs. ")
        
        pdf = TransformedPDF(title, workspace_name, orientation="P", unit="mm", format="A4")
        pdf.set_margins(15, 20, 15)
        pdf.set_auto_page_break(auto=True, margin=20)
        pdf.add_page()
        
        # Cover sheet / Section Header
        pdf.set_font("helvetica", "B", 18)
        pdf.set_text_color(11, 18, 32)
        pdf.cell(0, 10, title.upper(), border=0, align="L")
        pdf.ln(12)
        
        html_lines = []
        lines = text.split("\n")
        
        for line in lines:
            line_strip = line.strip()
            if not line_strip:
                html_lines.append("<br>")
                continue
                
            if line_strip.startswith("# "):
                html_lines.append(f"<h1>{line_strip[2:]}</h1>")
            elif line_strip.startswith("## "):
                html_lines.append(f"<h2>{line_strip[3:]}</h2>")
            elif line_strip.startswith("### "):
                html_lines.append(f"<h3>{line_strip[4:]}</h3>")
            elif line_strip.startswith("- ") or line_strip.startswith("* "):
                html_lines.append(f"<li>{line_strip[2:]}</li>")
            elif line_strip.startswith("|"):
                cells = [c.strip() for c in line_strip.split("|")[1:-1]]
                if not cells or "---" in line_strip:
                    continue
                html_lines.append("<tr>")
                for cell in cells:
                    if "Metric" in cell or "Value" in cell:
                        html_lines.append(f"<td><b>{cell}</b></td>")
                    else:
                        html_lines.append(f"<td>{cell}</td>")
                html_lines.append("</tr>")
            else:
                html_lines.append(f"<p>{line_strip}</p>")
                
        html_content = "".join(html_lines)
        html_content = re.sub(r"\*\*(.*?)\*\*", r"<b>\1</b>", html_content)
        html_content = re.sub(r"\*(.*?)\*", r"<i>\1</i>", html_content)
        
        # Enforce simple table wrapper if present
        if "<tr>" in html_content:
            html_content = html_content.replace("<tr>", "<table border='1' cellpadding='4'><tr>", 1)
            html_content += "</table>"
            
        pdf.set_font("helvetica", "", 10)
        pdf.write_html(html_content)
        pdf.output(out_path)
